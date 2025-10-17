import os
import json
import fitz  # PyMuPDF
# import cv2
import numpy as np
from PIL import Image
import requests
import base64
import tempfile
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple, Union, BinaryIO # MODIFIED: Added Any

# Custom exception for cancellation
class UploadCancelledError(Exception):
    """Custom exception to indicate that the upload was cancelled."""
    pass
from contextlib import contextmanager
from fastapi import HTTPException
from docx import Document  # python-docx để đọc file Word
from pptx import Presentation  # python-pptx để đọc file PowerPoint
import pandas as pd  # pandas để đọc file Excel
import uuid
import time
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type, RetryError
from requests.exceptions import RequestException, Timeout, ConnectionError

# Cấu hình logging chỉ xuất ra console
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

# Cấu hình Paddle OCR API
RETRY_MAX_ATTEMPTS = 3
RETRY_BACKOFF = 2

class OCRService:
    def __init__(self):
        self.supported_image_extensions = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif'}
        self.supported_pdf_extensions = {'.pdf'}
        self.supported_document_extensions = {'.docx', '.doc'}  # Word
        self.supported_presentation_extensions = {'.pptx', '.ppt'}  # PowerPoint
        self.supported_text_extensions = {'.txt'}  # Text files
        self.supported_excel_extensions = {'.xlsx', '.xls'}  # Excel files
        
        # Cấu hình Paddle OCR API từ .env file
        self.paddle_ocr_url = os.getenv("PADDLE_OCR_API_URL")
        if not self.paddle_ocr_url:
            raise ValueError("PADDLE_OCR_API_URL không được cấu hình trong file .env")
        
    @retry(
        stop=stop_after_attempt(RETRY_MAX_ATTEMPTS),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        retry=retry_if_exception_type((RequestException, Timeout, ConnectionError)),
        reraise=True
    )
    def _call_ocr_api_with_retry(self, files: dict, data: dict, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Gọi OCR API với cơ chế retry và kiểm tra hủy."""
        # --- MODIFICATION: CHECK FOR CANCELLATION ---
        if upload_id and document_service and not document_service.get_upload_info(upload_id):
            raise UploadCancelledError(f"Upload {upload_id} was cancelled before calling API.")
        # --- END MODIFICATION ---
        try:
            response = requests.post(
                self.paddle_ocr_url,
                files=files,
                data=data,
                timeout=30,
                verify=False  # Chỉ dùng trong môi trường dev, không dùng cho production
            )
            response.raise_for_status()
            return response.json()
            
        except requests.Timeout as e:
            print(f"⏱️  Timeout khi gọi OCR API: {str(e)}")
            raise
        except requests.ConnectionError as e:
            print(f"🔌 Lỗi kết nối đến OCR API: {str(e)}")
            raise
        except requests.RequestException as e:
            print(f"❌ Lỗi khi gọi OCR API: {str(e)}")
            raise

    def _process_ocr_result(self, api_result: Dict[str, Any]) -> Dict[str, Any]:
        """Xử lý kết quả OCR từ API"""
        try:
            ocr_results = api_result.get('result', [])
            full_text = []
            words_info = []
            
            # Xử lý format từ ocr-fullV2
            for block in ocr_results:
                if isinstance(block, list):
                    for result in block:
                        if isinstance(result, list) and len(result) >= 2:
                            # result[0] chứa bounding box, result[1] chứa (text, confidence)
                            bbox = result[0]
                            text_info = result[1]
                            
                            if isinstance(text_info, list) and len(text_info) >= 2:
                                text = text_info[0]
                                confidence = float(text_info[1]) * 100  # Chuyển về phần trăm
                                
                                full_text.append(text)
                                words_info.append({
                                    'text': text,
                                    'confidence': int(confidence),
                                    'bbox': {
                                        'x': int(bbox[0][0]),
                                        'y': int(bbox[0][1]),
                                        'width': int(bbox[2][0] - bbox[0][0]),
                                        'height': int(bbox[2][1] - bbox[0][1])
                                    }
                                })
            
            # In kết quả OCR ra console
            extracted_text = '\n'.join(full_text)
            
            return {
                'success': True,
                'text': extracted_text,
                'words': words_info,
                'total_words': len(words_info),
                'average_confidence': sum(w['confidence'] for w in words_info) / len(words_info) if words_info else 0
            }
            
        except Exception as e:
            print(f"❌ Lỗi khi xử lý kết quả OCR: {str(e)}")
            return {
                'success': False,
                'error': f'Lỗi xử lý kết quả OCR: {str(e)}',
                'text': '',
                'words': [],
                'total_words': 0,
                'average_confidence': 0
            }

    @contextmanager
    def _temp_image_file(self, image_data: Union[bytes, BinaryIO], suffix: str = '.png') -> str:
        """Tạo file tạm từ dữ liệu ảnh và tự động xóa sau khi sử dụng"""
        temp_file = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp:
                if isinstance(image_data, bytes):
                    temp.write(image_data)
                else:
                    temp.write(image_data.read())
                temp_file = temp.name
            yield temp_file
        finally:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.unlink(temp_file)
                except Exception as e:
                    logger.warning(f"Không thể xóa file tạm {temp_file}: {e}")

    def _ocr_image_from_bytes(self, image_data: Union[bytes, BinaryIO], file_name: str = 'image.png', upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Xử lý OCR từ dữ liệu ảnh dạng bytes hoặc file-like object"""
        with self._temp_image_file(image_data) as temp_file_path:
            return self._call_paddle_ocr_api(temp_file_path, file_name, upload_id=upload_id, document_service=document_service)

    def _call_paddle_ocr_api(self, image_path: str, file_name: str = None, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Gọi Paddle OCR API để trích xuất text từ hình ảnh với cơ chế retry"""
        if not file_name:
            file_name = os.path.basename(image_path)
            
        try:
            # Chuẩn bị file và form data cho API
            with open(image_path, "rb") as image_file:
                files = {
                    'file': (file_name, image_file, 'image/png')
                }
                data = {
                    'model': 'paddle',
                    'lang': 'vie'
                }
                
                # Gọi API với endpoint đã cấu hình
                logger.info(f"Gọi OCR API: {self.paddle_ocr_url}")
                logger.debug(f"Dữ liệu gửi đi: {data}")
                
                # Gọi API với retry và kiểm tra hủy
                api_result = self._call_ocr_api_with_retry(files, data, upload_id=upload_id, document_service=document_service)
                
                # Xử lý kết quả
                return self._process_ocr_result(api_result)
                
        except (RetryError, UploadCancelledError) as e: # MODIFIED: Catch UploadCancelledError
            if isinstance(e, UploadCancelledError):
                logger.info(str(e))
                raise # Re-raise to be caught by the calling function

            print(f"❌ Đã thử lại {RETRY_MAX_ATTEMPTS} lần nhưng vẫn lỗi: {str(e)}")
            return {
                'success': False,
                'error': f'OCR processing failed after {RETRY_MAX_ATTEMPTS} attempts: {str(e)}',
                'text': '',
                'words': [],
                'total_words': 0,
                'average_confidence': 0,
                'retry_count': RETRY_MAX_ATTEMPTS
            }
                
        except Exception as e:
            print(f"❌ Lỗi không xác định khi gọi OCR API: {str(e)}")
            return {
                'success': False,
                'error': f'Unexpected error during OCR processing: {str(e)}',
                'text': '',
                'words': [],
                'total_words': 0,
                'average_confidence': 0,
                'retry_count': 0 # MODIFIED: retry_count was not defined here
            }
        
    def extract_text_from_image(self, image_path: str, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Trích xuất text từ hình ảnh sử dụng Paddle OCR API"""
        try:
            print(f"Đang xử lý hình ảnh với Paddle OCR: {os.path.basename(image_path)}")
            
            # Kiểm tra file tồn tại
            if not os.path.exists(image_path):
                raise HTTPException(status_code=400, detail="File hình ảnh không tồn tại")
            
            # Gọi Paddle OCR API với tham số hủy
            result = self._call_paddle_ocr_api(image_path, upload_id=upload_id, document_service=document_service)
            
            # In chi tiết kết quả OCR
            if result.get('success', False):
                print(f"\n=== XỬ LÝ HÌNH ẢNH: {os.path.basename(image_path)} ===")
                print(f"Đã trích xuất thành công {result.get('total_words', 0)} từ")
                if result.get('text'):
                    print(f"Nội dung: {result.get('text')[:300]}{'...' if len(result.get('text', '')) > 300 else ''}")
                else:
                    print("CẢNH BÁO: Không có text nào được trích xuất từ hình ảnh!")
                print(f"=== KẾT THÚC XỬ LÝ ===\n")
            else:
                print(f"\n=== LỖI XỬ LÝ HÌNH ẢNH: {os.path.basename(image_path)} ===")
                print(f"Lỗi: {result.get('error', 'Không xác định')}")
                print(f"=== KẾT THÚC LỖI ===\n")
            
            return result
            
        except Exception as e:
            print(f"\n=== EXCEPTION XỬ LÝ HÌNH ẢNH: {os.path.basename(image_path) if image_path else 'Unknown'} ===")
            print(f"Exception: {str(e)}")
            print(f"=== KẾT THÚC EXCEPTION ===\n")
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'words': [],
                'total_words': 0,
                'average_confidence': 0
            }
    
    def _process_pdf_page(self, page, page_num: int, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Xử lý một trang PDF riêng lẻ"""
        try:
            # --- MODIFICATION: CHECK FOR CANCELLATION ---
            if upload_id and document_service and not document_service.get_upload_info(upload_id):
                raise UploadCancelledError(f"Upload {upload_id} was cancelled during PDF processing.")
            # --- END MODIFICATION ---

            page_text = page.get_text()
            images = page.get_images()
            has_images = len(images) > 0
            has_text = bool(page_text.strip())
            
            # Nếu không có hình ảnh hoặc không yêu cầu OCR, trả về text gốc
            if not has_images or not is_image:
                return {
                    'page_number': page_num + 1,
                    'text': page_text if has_text else '',
                    'has_images': has_images,
                    'has_text': has_text,
                    'success': True
                }
            
            # Nếu có hình ảnh và yêu cầu OCR
            try:
                # Tạo hình ảnh từ trang PDF
                pix = page.get_pixmap()
                with self._temp_image_file(pix.tobytes("png")) as temp_img_path:
                    # Thực hiện OCR với Paddle OCR API
                    ocr_result = self._call_paddle_ocr_api(temp_img_path, f"page_{page_num + 1}.png", upload_id=upload_id, document_service=document_service)
                    ocr_text = ocr_result.get('text', '') if ocr_result.get('success', False) else ''
                    
                    # Kết hợp text gốc và text từ OCR nếu cần
                    if has_text:
                        combined_text = f"{page_text}\n{ocr_text}"
                    else:
                        combined_text = ocr_text
                    
                    return {
                        'page_number': page_num + 1,
                        'text': combined_text,
                        'has_images': has_images,
                        'has_text': has_text or bool(ocr_text.strip()),
                        'ocr_raw_response': ocr_result,
                        'ocr_extracted_text': ocr_text,
                        'success': True
                    }
                    
            except Exception as ocr_error:
                logger.error(f"Lỗi OCR trang {page_num + 1}: {str(ocr_error)}")
                # Nếu có lỗi OCR nhưng có text gốc, vẫn trả về text gốc
                if has_text:
                    return {
                        'page_number': page_num + 1,
                        'text': page_text,
                        'has_images': has_images,
                        'has_text': True,
                        'ocr_error': str(ocr_error),
                        'success': True
                    }
                raise  # Ném lỗi nếu không có text gốc
                
        except Exception as e:
            logger.error(f"Lỗi khi xử lý trang {page_num + 1}: {str(e)}")
            return {
                'page_number': page_num + 1,
                'text': '',
                'has_images': False,
                'has_text': False,
                'error': str(e),
                'success': False
            }

    def extract_text_from_pdf(self, pdf_path: str, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Trích xuất text từ PDF với xử lý thông minh cho từng trang"""
        doc = None
        try:
            doc = fitz.open(pdf_path)
            full_text = []
            pages_info = []
            total_pages = len(doc)

            logger.info(f"Bắt đầu trích xuất text từ PDF '{os.path.basename(pdf_path)}' với {total_pages} trang...")

            # Xử lý từng trang
            for page_num in range(total_pages):
                try:
                    page = doc[page_num]
                    page_result = self._process_pdf_page(page, page_num, is_image, upload_id=upload_id, document_service=document_service)
                    
                    if page_result['success']:
                        full_text.append(page_result['text'])
                        pages_info.append(page_result)
                        
                        # Log tiến độ
                        if (page_num + 1) % 10 == 0 or (page_num + 1) == total_pages:
                            logger.info(f"Đã xử lý {page_num + 1}/{total_pages} trang")
                    else:
                        logger.warning(f"Không thể xử lý trang {page_num + 1}: {page_result.get('error', 'Lỗi không xác định')}")
                        pages_info.append(page_result)
                        
                except Exception as e:
                    logger.error(f"Lỗi không xử lý được trang {page_num + 1}: {str(e)}")
                    pages_info.append({
                        'page_number': page_num + 1,
                        'text': '',
                        'has_images': False,
                        'has_text': False,
                        'error': str(e),
                        'success': False
                    })

            # Tạo kết quả cuối cùng
            final_text = '\n\n'.join(filter(None, full_text))
            result = {
                'success': True,
                'text': final_text,
                'pages': pages_info,
                'total_pages': total_pages,
                'total_words': len(final_text.split()) if final_text else 0,
                'processed_pages': len([p for p in pages_info if p.get('success', False)])
            }
            
            logger.info(f"Hoàn thành xử lý PDF. Đã xử lý thành công {result['processed_pages']}/{total_pages} trang")
            return result
            
        except Exception as e:
            error_msg = f"Lỗi khi xử lý file PDF: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                'success': False,
                'error': error_msg,
                'text': '',
                'pages': [],
                'total_pages': 0,
                'total_words': 0,
                'processed_pages': 0
            }
            
        finally:
            if doc:
                doc.close()

    
    def extract_text_from_docx(self, file_path: str, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Trích xuất text từ file Word với xử lý thông minh text + OCR"""
        try:
            doc = Document(file_path)
            full_text = []
            sections_info = []

            logger.info(f"Bắt đầu xử lý file Word: {os.path.basename(file_path)}")

            # Xử lý paragraphs
            paragraph_texts = []
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    paragraph_texts.append(paragraph.text.strip())

            # Xử lý tables
            table_texts = []
            for table_idx, table in enumerate(doc.tables):
                table_content = []
                for row in table.rows:
                    row_content = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_content.append(cell.text.strip())
                    if row_content:
                        table_content.append(' | '.join(row_content))
                if table_content:
                    table_texts.append('\n'.join(table_content))

            # Xử lý images với OCR
            image_texts = []
            has_images = False
            try:
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        has_images = True
                        try:
                            # --- MODIFICATION: CHECK FOR CANCELLATION ---
                            if upload_id and document_service and not document_service.get_upload_info(upload_id):
                                raise UploadCancelledError(f"Upload {upload_id} was cancelled during DOCX processing.")
                            # --- END MODIFICATION ---

                            # Sử dụng _ocr_image_from_bytes thay vì xử lý file tạm
                            ocr_result = self._ocr_image_from_bytes(
                                rel.target_part.blob,
                                f"docx_image_{len(image_texts)}.png",
                                upload_id=upload_id, 
                                document_service=document_service
                            )
                            if ocr_result.get('success', False) and ocr_result.get('text'):
                                image_texts.append(ocr_result['text'])
                        except Exception as img_e:
                            logger.warning(f"Lỗi khi xử lý ảnh trong docx: {str(img_e)}")
                            continue
            except Exception as e:
                logger.error(f"Lỗi khi truy cập các mối quan hệ trong docx: {str(e)}", exc_info=True)

            # Kết hợp tất cả text
            all_texts = []
            if paragraph_texts:
                all_texts.extend(paragraph_texts)
            if table_texts:
                all_texts.extend(table_texts)
            if image_texts:
                all_texts.extend(image_texts)

            final_text = '\n\n'.join(all_texts)

            result = {
                'success': True,
                'text': final_text,
                'total_words': len(final_text.split()),
                'sections': {
                    'paragraphs': len(paragraph_texts),
                    'tables': len(table_texts),
                    'images': len(image_texts)
                },
                'has_images': has_images,
                'processing_summary': {
                    'paragraph_text_extracted': len(paragraph_texts) > 0,
                    'table_text_extracted': len(table_texts) > 0,
                    'image_ocr_performed': len(image_texts) > 0
                }
            }

            logger.info(f"Hoàn thành xử lý file Word. Tổng số từ: {result['total_words']}")
            return result

        except Exception as e:
            error_msg = f"Lỗi khi xử lý file Word: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                'success': False,
                'error': error_msg,
                'text': '',
                'total_words': 0
            }
    
    def extract_text_from_pptx(self, file_path: str, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Trích xuất text từ file PowerPoint với xử lý thông minh text + OCR"""
        try:
            logger.info(f"Bắt đầu xử lý file PowerPoint: {os.path.basename(file_path)}")
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            all_slide_texts = []
            slides_info = []
            
            logger.info(f"Tổng số slide cần xử lý: {total_slides}")
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    # --- MODIFICATION: CHECK FOR CANCELLATION ---
                    if upload_id and document_service and not document_service.get_upload_info(upload_id):
                        raise UploadCancelledError(f"Upload {upload_id} was cancelled during PPTX processing.")
                    # --- END MODIFICATION ---

                    slide_text_parts = []
                    slide_image_texts = []
                    has_images = False
                    
                    # Xử lý text shapes
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text_parts.append(shape.text.strip())
                            
                            # Kiểm tra và xử lý images nếu được yêu cầu
                            if is_image and hasattr(shape, 'image') and hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture type
                                has_images = True
                                try:
                                    # Sử dụng _ocr_image_from_bytes để xử lý ảnh
                                    ocr_result = self._ocr_image_from_bytes(
                                        shape.image.blob,
                                        f"pptx_slide_{slide_idx + 1}_img_{len(slide_image_texts)}.png",
                                        upload_id=upload_id, 
                                        document_service=document_service
                                    )
                                    if ocr_result.get('success', False) and ocr_result.get('text'):
                                        slide_image_texts.append(ocr_result['text'])
                                        logger.debug(f"Đã xử lý ảnh {len(slide_image_texts)} trong slide {slide_idx + 1}")
                                        
                                except Exception as img_e:
                                    logger.warning(f"Lỗi khi xử lý ảnh trong slide {slide_idx + 1}: {str(img_e)}")
                                    continue
                                    
                        except Exception as shape_e:
                            logger.warning(f"Lỗi khi xử lý shape trong slide {slide_idx + 1}: {str(shape_e)}")
                            continue
                    
                    # Kết hợp text và OCR cho slide này
                    slide_all_text = []
                    if slide_text_parts:
                        slide_all_text.extend(slide_text_parts)
                    if slide_image_texts:
                        slide_all_text.append("=== NỘI DUNG TỪ HÌNH ẢNH ===")
                        slide_all_text.extend(slide_image_texts)
                    
                    slide_final_text = '\n'.join(slide_all_text)
                    
                    # Log tiến độ
                    if (slide_idx + 1) % 5 == 0 or (slide_idx + 1) == total_slides:
                        logger.info(f"Đã xử lý {slide_idx + 1}/{total_slides} slide")
                    
                    # Lưu thông tin slide
                    slide_info = {
                        'slide_number': slide_idx + 1,
                        'text': slide_final_text,
                        'has_text': len(slide_text_parts) > 0,
                        'has_images': has_images,
                        'text_shapes_count': len(slide_text_parts),
                        'images_ocr_count': len(slide_image_texts),
                        'success': True
                    }
                    
                    slides_info.append(slide_info)
                    
                    if slide_final_text:
                        all_slide_texts.append(slide_final_text)
                        
                except Exception as slide_e:
                    error_msg = f"Lỗi khi xử lý slide {slide_idx + 1}: {str(slide_e)}"
                    logger.error(error_msg, exc_info=True)
                    slides_info.append({
                        'slide_number': slide_idx + 1,
                        'error': error_msg,
                        'has_text': False,
                        'has_images': False,
                        'success': False
                    })
            
            # Tạo kết quả cuối cùng
            final_text = '\n\n'.join(all_slide_texts)
            
            result = {
                'success': True,
                'text': final_text,
                'total_slides': total_slides,
                'total_words': len(final_text.split()),
                'slides': slides_info,
                'processing_summary': {
                    'slides_processed': len(slides_info),
                    'slides_with_text': len([s for s in slides_info if s.get('has_text', False)]),
                    'slides_with_images': len([s for s in slides_info if s.get('has_images', False)]),
                    'total_images_processed': sum(s.get('images_ocr_count', 0) for s in slides_info)
                }
            }
            
            logger.info(
                f"Hoàn thành xử lý PowerPoint. "
                f"Tổng số slide: {result['total_slides']}, "
                f"Tổng số từ: {result['total_words']}, "
                f"Số ảnh đã xử lý: {result['processing_summary']['total_images_processed']}"
            )
            
            return result
            
        except Exception as e:
            error_msg = f"Lỗi khi xử lý file PowerPoint: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return {
                'success': False,
                'error': error_msg,
                'text': '',
                'total_slides': 0,
                'total_words': 0,
                'slides': []
            }

    
    def extract_text_from_txt(self, file_path: str) -> Dict[str, Any]:
        """Đọc nội dung từ file text"""
        try:
            print("Đang đọc file text...")
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            return {
                'success': True,
                'text': text,
                'total_words': len(text.split())
            }
            
        except UnicodeDecodeError:
            # Thử lại với encoding khác nếu utf-8 không đọc được
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    text = f.read()
                return {
                    'success': True,
                    'text': text,
                    'total_words': len(text.split())
                }
            except Exception as e:
                return {
                    'success': False,
                    'error': str(e),
                    'text': '',
                    'total_words': 0
                }
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'total_words': 0
            }
    
    def extract_text_from_excel(self, file_path: str, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Trích xuất text từ file Excel với xử lý thông minh text + OCR"""
        try:
            # Đọc tất cả các sheet
            excel_file = pd.ExcelFile(file_path)
            all_sheets_text = []
            sheets_info = []
            
            for sheet_name in excel_file.sheet_names:
                # --- MODIFICATION: CHECK FOR CANCELLATION ---
                if upload_id and document_service and not document_service.get_upload_info(upload_id):
                    raise UploadCancelledError(f"Upload {upload_id} was cancelled during Excel processing.")
                # --- END MODIFICATION ---

                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Xử lý text data từ cells
                sheet_text_parts = []
                
                # Thêm tên sheet
                sheet_text_parts.append(f"=== SHEET: {sheet_name} ===")
                
                # Xử lý headers
                if not df.empty:
                    headers = [str(col) for col in df.columns if str(col) != 'nan']
                    if headers:
                        sheet_text_parts.append("Headers: " + " | ".join(headers))
                
                # Xử lý data rows
                for index, row in df.iterrows():
                    row_data = []
                    for col in df.columns:
                        cell_value = row[col]
                        if pd.notna(cell_value) and str(cell_value).strip():
                            row_data.append(str(cell_value).strip())
                    
                    if row_data:
                        sheet_text_parts.append(" | ".join(row_data))
                
                # Xử lý images trong Excel (nếu có)
                sheet_image_texts = []
                has_images = False
                
                try:
                    # Kiểm tra images trong workbook
                    from openpyxl import load_workbook
                    wb = load_workbook(file_path)
                    
                    if sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        
                        # Tìm images trong sheet
                        if hasattr(ws, '_images') and ws._images:
                            has_images = True
                            
                            for img_idx, img in enumerate(ws._images):
                                try:
                                    # Lưu image tạm thời
                                    temp_img_path = f"temp_excel_{sheet_name}_{img_idx}.png"
                                    
                                    # Lưu image data
                                    with open(temp_img_path, "wb") as f:
                                        f.write(img._data())
                                    
                                    # OCR image
                                    ocr_result = self._call_paddle_ocr_api(temp_img_path, upload_id=upload_id, document_service=document_service)
                                    if ocr_result.get('success', False) and ocr_result.get('text'):
                                        sheet_image_texts.append(ocr_result['text'])
                                    
                                    # Xóa file tạm
                                    if os.path.exists(temp_img_path):
                                        os.remove(temp_img_path)
                                        
                                except Exception:
                                    continue
                                    
                except Exception:
                    # Nếu không thể xử lý images, tiếp tục với text
                    pass
                
                # Kết hợp text và OCR cho sheet này
                sheet_all_text = []
                if sheet_text_parts:
                    sheet_all_text.extend(sheet_text_parts)
                if sheet_image_texts:
                    sheet_all_text.append("=== OCR FROM IMAGES ===")
                    sheet_all_text.extend(sheet_image_texts)
                
                sheet_final_text = '\n'.join(sheet_all_text)
                if sheet_final_text:
                    all_sheets_text.append(sheet_final_text)
                
                sheets_info.append({
                    'sheet_name': sheet_name,
                    'text': sheet_final_text,
                    'rows_count': len(df),
                    'columns_count': len(df.columns),
                    'has_images': has_images,
                    'images_ocr_count': len(sheet_image_texts),
                    'has_data': not df.empty
                })
            
            final_text = '\n\n'.join(all_sheets_text)
            
            return {
                'success': True,
                'text': final_text,
                'total_words': len(final_text.split()),
                'sheets': sheets_info,
                'processing_summary': {
                    'total_sheets': len(excel_file.sheet_names),
                    'sheets_with_data': len([s for s in sheets_info if s['has_data']]),
                    'sheets_with_images': len([s for s in sheets_info if s['has_images']]),
                    'total_images_processed': sum(s['images_ocr_count'] for s in sheets_info)
                }
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'sheets': [],
                'total_words': 0
            }

    def _extract_text_from_excel_images(self, file_path: str) -> Dict[str, Any]:
        """Trích xuất text từ hình ảnh trong file Excel bằng OCR"""
        try:
            print("Đang xử lý hình ảnh trong file Excel...")
            workbook = load_workbook(file_path)
            full_text = []
            sheets_data = []
            image_count = 0
            temp_dir = Path("temp_excel_images")
            temp_dir.mkdir(exist_ok=True)

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                sheet_images = []

                # Lấy text thông thường trước
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        sheet_text.extend(row_text)

                # Xử lý hình ảnh trong sheet (nếu có)
                try:
                    # Kiểm tra xem sheet có hình ảnh không
                    if hasattr(sheet, '_images') and sheet._images:
                        for img in sheet._images:
                            try:
                                image_count += 1
                                temp_path = temp_dir / f"excel_image_{image_count}.png"
                                
                                # Lưu hình ảnh
                                with open(temp_path, "wb") as f:
                                    f.write(img._data())

                                # OCR hình ảnh với Paddle OCR API
                                ocr_result = self._call_paddle_ocr_api(str(temp_path))
                                if ocr_result['success'] and ocr_result['text']:
                                    print(f"\n=== OCR HÌNH ẢNH EXCEL #{image_count} ===")
                                    print(f"Sheet: {sheet_name}")
                                    print(f"Nội dung: {ocr_result['text'][:200]}{'...' if len(ocr_result['text']) > 200 else ''}")
                                    print(f"=== KẾT THÚC HÌNH ẢNH #{image_count} ===\n")
                                    sheet_images.append({
                                        'image_id': image_count,
                                        'text': ocr_result['text']
                                    })
                                    sheet_text.append(ocr_result['text'])

                            except Exception as e:
                                print(f"Lỗi khi xử lý ảnh {image_count}: {str(e)}")
                except Exception as e:
                    print(f"Không thể truy cập hình ảnh trong sheet {sheet_name}: {str(e)}")

                # Thêm kết quả của sheet vào danh sách
                if sheet_text:
                    sheets_data.append({
                        'sheet_name': sheet_name,
                        'text': '\n'.join(sheet_text),
                        'images': sheet_images
                    })
                    full_text.append(f"Sheet: {sheet_name}")
                    full_text.extend(sheet_text)

            # Xóa thư mục tạm nếu tồn tại
            if temp_dir.exists():
                for item in temp_dir.iterdir():
                    item.unlink()
                temp_dir.rmdir()

            text = '\n\n'.join(full_text)
            return {
                'success': True,
                'text': text,
                'sheets': sheets_data,
                'total_sheets': len(sheets_data),
                'total_images': image_count,
                'total_words': len(text.split())
            }

        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'text': '',
                'sheets': [],
                'total_sheets': 0,
                'total_images': 0,
                'total_words': 0
            }

    def process_file(self, file_path: str, is_image: bool = False, upload_id: Optional[str] = None, document_service: Optional[Any] = None) -> Dict[str, Any]:
        """Xử lý file dựa trên extension và trả về kết quả OCR"""
        # --- MODIFICATION: CHECK FOR CANCELLATION AT THE START ---
        if upload_id and document_service and not document_service.get_upload_info(upload_id):
            raise UploadCancelledError(f"Upload {upload_id} was cancelled before processing.")
        # --- END MODIFICATION ---

        file_path = Path(file_path)
        
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="File không tồn tại")
        
        file_extension = file_path.suffix.lower()
        
        try:
            # Xử lý PDF
            if file_extension in self.supported_pdf_extensions:
                result = self.extract_text_from_pdf(str(file_path), is_image, upload_id=upload_id, document_service=document_service)
                result['file_type'] = 'pdf'
            
            # Xử lý hình ảnh
            elif file_extension in self.supported_image_extensions:
                result = self.extract_text_from_image(str(file_path), upload_id=upload_id, document_service=document_service)
                result['file_type'] = 'image'
            
            # Xử lý Word
            elif file_extension in self.supported_document_extensions:
                result = self.extract_text_from_docx(str(file_path), is_image, upload_id=upload_id, document_service=document_service)
                result['file_type'] = 'word'
            
            # Xử lý PowerPoint
            elif file_extension in self.supported_presentation_extensions:
                result = self.extract_text_from_pptx(str(file_path), is_image, upload_id=upload_id, document_service=document_service)
                result['file_type'] = 'powerpoint'
            
            # Xử lý Excel
            elif file_extension in self.supported_excel_extensions:
                result = self.extract_text_from_excel(str(file_path), is_image, upload_id=upload_id, document_service=document_service)
                result['file_type'] = 'excel'
            
            # Xử lý Text
            elif file_extension in self.supported_text_extensions:
                result = self.extract_text_from_txt(str(file_path))
                result['file_type'] = 'text'
            
            else:
                supported_extensions = (
                    self.supported_pdf_extensions |
                    self.supported_image_extensions |
                    self.supported_document_extensions |
                    self.supported_presentation_extensions |
                    self.supported_excel_extensions |
                    self.supported_text_extensions
                )
                raise HTTPException(
                    status_code=400,
                    detail=f"Định dạng file không được hỗ trợ. Hỗ trợ: {', '.join(supported_extensions)}"
                )
            
            result['file_name'] = file_path.name
            return result
            
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Lỗi xử lý file: {str(e)}")

# Khởi tạo service
ocr_service = OCRService()